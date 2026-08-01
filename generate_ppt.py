"""
========================================================================================
PowerPoint (PPTX) Generator for ME4V33 PLM Digital Twin Project
Course: ME4V33 – Product Lifecycle Management
Mentor: Dhanesh Babu
Team: Thirumalairajan U & Tharunkumar RP
========================================================================================
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # 16:9 Widescreen aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette Definitions
    DARK_BG = RGBColor(15, 23, 42)       # #0F172A Dark Slate
    CARD_BG = RGBColor(30, 41, 59)       # #1E293B Card Background
    CYAN_TITLE = RGBColor(56, 189, 248)   # #38BDF8 Electric Cyan
    WHITE_TEXT = RGBColor(248, 250, 252) # #F8FAFC
    MUTED_TEXT = RGBColor(148, 163, 184) # #94A3B8
    GREEN_ACCENT = RGBColor(16, 185, 129)# #10B981
    YELLOW_ACCENT = RGBColor(245, 158, 11)# #F59E0B
    RED_ACCENT = RGBColor(239, 68, 68)   # #EF4444

    blank_slide_layout = prs.slide_layouts[6]

    def set_slide_background(slide, color=DARK_BG):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, title_text, category_text="ME4506 – PRODUCT LIFECYCLE MANAGEMENT"):
        # Category Subtitle
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = CYAN_TITLE
        
        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = WHITE_TEXT

    # ----------------------------------------------------------------------------------
    # SLIDE 1: TITLE SLIDE
    # ----------------------------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide1)

    # Accent decorative box
    shape1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(0.15), Inches(5.0))
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = CYAN_TITLE
    shape1.line.color.rgb = CYAN_TITLE

    # Main Title
    tb1 = slide1.shapes.add_textbox(Inches(1.2), Inches(1.2), Inches(11.2), Inches(2.2))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "Digital Twin-Enabled AI Framework for Predictive Maintenance"
    p1.font.size = Pt(32)
    p1.font.bold = True
    p1.font.color.rgb = CYAN_TITLE

    p1_sub = tf1.add_paragraph()
    p1_sub.text = "Product Lifecycle Optimization of Heavy Equipment Engine Components"
    p1_sub.font.size = Pt(20)
    p1_sub.font.bold = True
    p1_sub.font.color.rgb = WHITE_TEXT

    # Course Metadata Card
    card1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(3.8), Inches(11.0), Inches(2.4))
    card1.fill.solid()
    card1.fill.fore_color.rgb = CARD_BG
    card1.line.color.rgb = CYAN_TITLE

    tf_c1 = card1.text_frame
    tf_c1.word_wrap = True
    
    p_meta1 = tf_c1.paragraphs[0]
    p_meta1.text = "📚 COURSE: ME4506 – Product Lifecycle Management (PLM)"
    p_meta1.font.size = Pt(16)
    p_meta1.font.bold = True
    p_meta1.font.color.rgb = WHITE_TEXT

    p_meta2 = tf_c1.add_paragraph()
    p_meta2.text = "👨‍🏫 FACULTY MENTOR: Dhanesh Babu"
    p_meta2.font.size = Pt(16)
    p_meta2.font.bold = True
    p_meta2.font.color.rgb = WHITE_TEXT

    p_meta3 = tf_c1.add_paragraph()
    p_meta3.text = "👥 PROJECT TEAM: Thirumalairajan U  &  Tharunkumar RP"
    p_meta3.font.size = Pt(16)
    p_meta3.font.bold = True
    p_meta3.font.color.rgb = GREEN_ACCENT

    # ----------------------------------------------------------------------------------
    # SLIDE 2: PROJECT ABSTRACT & PROBLEM STATEMENT
    # ----------------------------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide2)
    add_header(slide2, "Project Abstract & Problem Context")

    # Left Column: Problem Statement
    box_p = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0))
    box_p.fill.solid()
    box_p.fill.fore_color.rgb = CARD_BG
    box_p.line.color.rgb = RED_ACCENT
    tf_p = box_p.text_frame
    tf_p.word_wrap = True
    p_ph = tf_p.paragraphs[0]
    p_ph.text = "⚠️ Limitations of Conventional Maintenance"
    p_ph.font.size = Pt(18)
    p_ph.font.bold = True
    p_ph.font.color.rgb = RED_ACCENT

    p_pb1 = tf_p.add_paragraph()
    p_pb1.text = "• Extreme Engine Stress: High mechanical loads, continuous vibration, and thermal pulses accelerate bearing and piston wear."
    p_pb1.font.size = Pt(14)
    p_pb1.font.color.rgb = MUTED_TEXT

    p_pb2 = tf_p.add_paragraph()
    p_pb2.text = "• Reactive Maintenance: Fixes components after catastrophic failure, incurring average $48,000+ repair costs and unplanned downtime."
    p_pb2.font.size = Pt(14)
    p_pb2.font.color.rgb = MUTED_TEXT

    p_pb3 = tf_p.add_paragraph()
    p_pb3.text = "• Preventive Maintenance: Fixed-schedule servicing leads to premature replacement of healthy components."
    p_pb3.font.size = Pt(14)
    p_pb3.font.color.rgb = MUTED_TEXT

    # Right Column: Proposed AI Digital Twin Solution
    box_s = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    box_s.fill.solid()
    box_s.fill.fore_color.rgb = CARD_BG
    box_s.line.color.rgb = GREEN_ACCENT
    tf_s = box_s.text_frame
    tf_s.word_wrap = True
    p_sh = tf_s.paragraphs[0]
    p_sh.text = "💡 Proposed AI + Digital Twin Framework"
    p_sh.font.size = Pt(18)
    p_sh.font.bold = True
    p_sh.font.color.rgb = GREEN_ACCENT

    p_sb1 = tf_s.add_paragraph()
    p_sb1.text = "• Digital Twin Virtualization: Real-time telemetry monitoring (Temp, Vibration, Oil Pressure, RPM, Load, Duty Hours)."
    p_sb1.font.size = Pt(14)
    p_sb1.font.color.rgb = MUTED_TEXT

    p_sb2 = tf_s.add_paragraph()
    p_sb2.text = "• Parallel Machine Learning: Random Forest Classifier (96.67% accuracy) & RUL Regressor (R² = 0.97)."
    p_sb2.font.size = Pt(14)
    p_sb2.font.color.rgb = MUTED_TEXT

    p_sb3 = tf_s.add_paragraph()
    p_sb3.text = "• PLM Decision Layer: Converts AI predictions into 3 actionable lifecycle decisions (Operate, Maintain, Replace)."
    p_sb3.font.size = Pt(14)
    p_sb3.font.color.rgb = MUTED_TEXT

    # ----------------------------------------------------------------------------------
    # SLIDE 3: KEY OBJECTIVES & PROJECT SCOPE
    # ----------------------------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide3)
    add_header(slide3, "Project Key Objectives & Scope")

    objs = [
        ("Objective 1: Dataset Engineering", "Build a synthetic/real sensor dataset (1,200 records) reflecting realistic engine degradation physics and operational stress parameters.", CYAN_TITLE),
        ("Objective 2: Machine Learning Models", "Train a Random Forest classifier for health state prediction (Healthy, Warning, Failure) and a regressor for Remaining Useful Life (RUL) estimation.", GREEN_ACCENT),
        ("Objective 3: Digital Twin Dashboard", "Develop an interactive Streamlit digital twin dashboard to mirror physical component state, displaying RUL gauges, ISO 10816 vibration classes, and CAD wear maps.", YELLOW_ACCENT),
        ("Objective 4: PLM Decision Integration", "Integrate AI outputs directly into a PLM decision layer to recommend Operate, Maintain, or Replace actions with financial ROI cost optimization.", RED_ACCENT)
    ]

    top_pos = 1.8
    for title, desc, col in objs:
        obox = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(top_pos), Inches(11.7), Inches(1.1))
        obox.fill.solid()
        obox.fill.fore_color.rgb = CARD_BG
        obox.line.color.rgb = col
        
        tf_o = obox.text_frame
        tf_o.word_wrap = True
        p_oh = tf_o.paragraphs[0]
        p_oh.text = title
        p_oh.font.size = Pt(16)
        p_oh.font.bold = True
        p_oh.font.color.rgb = col

        p_od = tf_o.add_paragraph()
        p_od.text = desc
        p_od.font.size = Pt(13)
        p_od.font.color.rgb = WHITE_TEXT
        
        top_pos += 1.35

    # ----------------------------------------------------------------------------------
    # SLIDE 4: SYSTEM ARCHITECTURE & DIGITAL THREAD PIPELINE
    # ----------------------------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide4)
    add_header(slide4, "Digital Thread Telemetry Pipeline Architecture")

    steps = [
        ("1. Edge Telemetry", "Temp, Vibration, Oil Press, RPM, Load, Hours", CYAN_TITLE),
        ("2. Data Preprocessing", "StandardScaler & LabelEncoder (Python)", GREEN_ACCENT),
        ("3. ML Inference Engine", "Random Forest Classifier & RUL Regressor", YELLOW_ACCENT),
        ("4. Digital Twin Dashboard", "Real-Time Telemetry & ISO 10816 Diagnostics", CYAN_TITLE),
        ("5. Enterprise PLM", "PTC Windchill sBOM Export & ECR Actions", RED_ACCENT)
    ]

    left_pos = 0.8
    for stitle, sdesc, scol in steps:
        sbox = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left_pos), Inches(2.2), Inches(2.1), Inches(4.2))
        sbox.fill.solid()
        sbox.fill.fore_color.rgb = CARD_BG
        sbox.line.color.rgb = scol
        
        tf_st = sbox.text_frame
        tf_st.word_wrap = True
        p_sth = tf_st.paragraphs[0]
        p_sth.text = stitle
        p_sth.font.size = Pt(15)
        p_sth.font.bold = True
        p_sth.font.color.rgb = scol

        p_std = tf_st.add_paragraph()
        p_std.text = f"\n{sdesc}"
        p_std.font.size = Pt(13)
        p_std.font.color.rgb = WHITE_TEXT

        left_pos += 2.4

    # ----------------------------------------------------------------------------------
    # SLIDE 5: MACHINE LEARNING EVALUATION RESULTS
    # ----------------------------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide5)
    add_header(slide5, "Machine Learning Performance & Validation")

    # Left: Classification Results
    c_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0))
    c_box.fill.solid()
    c_box.fill.fore_color.rgb = CARD_BG
    c_box.line.color.rgb = GREEN_ACCENT
    tf_cb = c_box.text_frame
    tf_cb.word_wrap = True
    
    p_c1 = tf_cb.paragraphs[0]
    p_c1.text = "🎯 Classifier Results (Random Forest)"
    p_c1.font.size = Pt(18)
    p_c1.font.bold = True
    p_c1.font.color.rgb = GREEN_ACCENT

    p_c2 = tf_cb.add_paragraph()
    p_c2.text = "• Accuracy: 96.67%\n• Weighted Precision: 96.71%\n• Weighted Recall: 96.67%\n• F1 Score: 96.68%\n\nConfusion Matrix:"
    p_c2.font.size = Pt(14)
    p_c2.font.color.rgb = WHITE_TEXT

    p_cm = tf_cb.add_paragraph()
    p_cm.text = "              Failure  Healthy  Warning\nFailure       67         0        3\nHealthy        0        68        2\nWarning        3         0       97"
    p_cm.font.size = Pt(12)
    p_cm.font.bold = True
    p_cm.font.color.rgb = CYAN_TITLE

    # Right: Regression Results
    r_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    r_box.fill.solid()
    r_box.fill.fore_color.rgb = CARD_BG
    r_box.line.color.rgb = CYAN_TITLE
    tf_rb = r_box.text_frame
    tf_rb.word_wrap = True
    
    p_r1 = tf_rb.paragraphs[0]
    p_r1.text = "📈 RUL Regressor Results (Random Forest)"
    p_r1.font.size = Pt(18)
    p_r1.font.bold = True
    p_r1.font.color.rgb = CYAN_TITLE

    p_r2 = tf_rb.add_paragraph()
    p_r2.text = "• Mean Absolute Error (MAE): 42.40 hours\n• Root Mean Sq Error (RMSE): 54.80 hours\n• R² Score: 0.9717 (97.17% Variance Explained)\n\nKey Insights:\n• High precision RUL estimation enables proactive maintenance scheduling 150+ hours prior to catastrophic failure hazard."
    p_r2.font.size = Pt(14)
    p_r2.font.color.rgb = WHITE_TEXT

    # ----------------------------------------------------------------------------------
    # SLIDE 6: ISO STANDARDS & ENGINEERING DIAGNOSTICS
    # ----------------------------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide6)
    add_header(slide6, "ISO Engineering Standards Compliance")

    iso_items = [
        ("ISO 10816 / ISO 20816", "Mechanical Vibration Severity", "Monitors non-rotating bearing vibration velocity (mm/s):\n• Class A: Good / Excellent (< 2.8 mm/s)\n• Class B: Acceptable Operation (2.8 - 7.1 mm/s)\n• Class C: Unsatisfactory (7.1 - 11.2 mm/s)\n• Class D: Damage Hazard (> 11.2 mm/s)", GREEN_ACCENT),
        ("ISO 4406", "Fluid Cleanliness Rating", "Quantifies solid particle contamination per mL across 3 size thresholds (≥4µm / ≥6µm / ≥14µm). High particle counts correlate with bearing spalling and cylinder scuffing.", CYAN_TITLE),
        ("ISO 13374 & ISO 23247", "Digital Twin Standards", "Standardized Condition Monitoring architecture (Data Acquisition → State Detection → Health Assessment → Prognostics RUL → Advisory Generation).", YELLOW_ACCENT)
    ]

    top_iso = 1.8
    for is1, is2, is3, icol in iso_items:
        ibox = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(top_iso), Inches(11.7), Inches(1.5))
        ibox.fill.solid()
        ibox.fill.fore_color.rgb = CARD_BG
        ibox.line.color.rgb = icol
        
        tf_i = ibox.text_frame
        tf_i.word_wrap = True
        p_ih = tf_i.paragraphs[0]
        p_ih.text = f"{is1} – {is2}"
        p_ih.font.size = Pt(16)
        p_ih.font.bold = True
        p_ih.font.color.rgb = icol

        p_id = tf_i.add_paragraph()
        p_id.text = is3
        p_id.font.size = Pt(13)
        p_id.font.color.rgb = WHITE_TEXT
        
        top_iso += 1.75

    # ----------------------------------------------------------------------------------
    # SLIDE 7: DIGITAL TWIN DASHBOARD & SUBSYSTEM MONITORS
    # --------------------------------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide7)
    add_header(slide7, "Digital Twin Telemetry Dashboard Features")

    feats = [
        ("🌐 Live Telemetry", "Real-time health status badges (Healthy 🟢, Warning 🟡, Failure Risk 🔴) & confidence scores."),
        ("📊 RUL Gauge Indicator", "Visual Plotly gauge tracking remaining operational hours against overhaul thresholds."),
        ("🧩 Subsystem Wear Map", "Dynamic wear index computation for Crankshaft Bearings, Pistons, Oil Pump, and Turbocharger."),
        ("🎛️ Simulation Modes", "Pre-configured operational stress tests (Thermal Overload, Vibration Fault, Lubrication Starvation).")
    ]

    left_f = 0.8
    for fhead, fbody in feats:
        fbox = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left_f), Inches(1.8), Inches(2.7), Inches(5.0))
        fbox.fill.solid()
        fbox.fill.fore_color.rgb = CARD_BG
        fbox.line.color.rgb = CYAN_TITLE
        
        tf_f = fbox.text_frame
        tf_f.word_wrap = True
        p_fh = tf_f.paragraphs[0]
        p_fh.text = fhead
        p_fh.font.size = Pt(16)
        p_fh.font.bold = True
        p_fh.font.color.rgb = CYAN_TITLE

        p_fb = tf_f.add_paragraph()
        p_fb.text = f"\n{fbody}"
        p_fb.font.size = Pt(13)
        p_fb.font.color.rgb = WHITE_TEXT

        left_f += 2.95

    # ----------------------------------------------------------------------------------
    # SLIDE 8: EXPLAINABLE AI (XAI) & GENERATIVE AI COPILOT
    # ----------------------------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide8)
    add_header(slide8, "Explainable AI (XAI) & Generative AI Copilot")

    # Left: XAI
    x_box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0))
    x_box.fill.solid()
    x_box.fill.fore_color.rgb = CARD_BG
    x_box.line.color.rgb = CYAN_TITLE
    tf_xb = x_box.text_frame
    tf_xb.word_wrap = True
    p_x1 = tf_xb.paragraphs[0]
    p_x1.text = "🔍 Feature Importance (XAI)"
    p_x1.font.size = Pt(18)
    p_x1.font.bold = True
    p_x1.font.color.rgb = CYAN_TITLE

    p_x2 = tf_xb.add_paragraph()
    p_x2.text = "• Quantifies mathematical feature weights to identify principal degradation drivers.\n• Identifies primary failure mechanisms (e.g. vibration velocity or oil pressure drop).\n• Generates automated plain-English diagnostic evaluations for engineers."
    p_x2.font.size = Pt(14)
    p_x2.font.color.rgb = WHITE_TEXT

    # Right: Generative AI
    g_box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    g_box.fill.solid()
    g_box.fill.fore_color.rgb = CARD_BG
    g_box.line.color.rgb = GREEN_ACCENT
    tf_gb = g_box.text_frame
    tf_gb.word_wrap = True
    p_g1 = tf_gb.paragraphs[0]
    p_g1.text = "🤖 Generative AI Copilot & RCFA"
    p_g1.font.size = Pt(18)
    p_g1.font.bold = True
    p_g1.font.color.rgb = GREEN_ACCENT

    p_g2 = tf_gb.add_paragraph()
    p_g2.text = "• Interactive AI Chatbot: Answers custom technical questions on bearing repair & PLM strategy.\n• Root Cause Failure Analysis (RCFA): Recommends exact overhaul kits (PM-250, PM-1000, PM-3000).\n• Google Gemini API Support with intelligent offline engineering fallback."
    p_g2.font.size = Pt(14)
    p_g2.font.color.rgb = WHITE_TEXT

    # ----------------------------------------------------------------------------------
    # SLIDE 9: PLM DECISION LAYER & FINANCIAL ROI
    # ----------------------------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide9)
    add_header(slide9, "PLM Decision Layer & Financial ROI Optimizer")

    # Top: 3 Lifecycle Actions
    actions = [
        ("1. Operate (Healthy 🟢)", "Continue standard operation & PM-250 servicing.", GREEN_ACCENT),
        ("2. Maintain (Warning 🟡)", "Schedule inspection within 150 hrs & PM-1000 overhaul.", YELLOW_ACCENT),
        ("3. Replace (Failure 🔴)", "Initiate replacement planning & PM-3000 emergency kit.", RED_ACCENT)
    ]

    left_a = 0.8
    for act1, act2, acol in actions:
        abox = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left_a), Inches(1.8), Inches(3.7), Inches(1.8))
        abox.fill.solid()
        abox.fill.fore_color.rgb = CARD_BG
        abox.line.color.rgb = acol
        tf_a = abox.text_frame
        tf_a.word_wrap = True
        p_act1 = tf_a.paragraphs[0]
        p_act1.text = act1
        p_act1.font.size = Pt(15)
        p_act1.font.bold = True
        p_act1.font.color.rgb = acol
        p_act2 = tf_a.add_paragraph()
        p_act2.text = act2
        p_act2.font.size = Pt(12)
        p_act2.font.color.rgb = WHITE_TEXT
        left_a += 3.95

    # Bottom: Financial ROI & PTC Windchill
    f_box = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.9), Inches(11.7), Inches(2.9))
    f_box.fill.solid()
    f_box.fill.fore_color.rgb = CARD_BG
    f_box.line.color.rgb = CYAN_TITLE
    tf_f = f_box.text_frame
    tf_f.word_wrap = True
    p_fh = tf_f.paragraphs[0]
    p_fh.text = "💰 Financial ROI & Enterprise PLM Integration (PTC Windchill)"
    p_fh.font.size = Pt(17)
    p_fh.font.bold = True
    p_fh.font.color.rgb = CYAN_TITLE

    p_fb = tf_f.add_paragraph()
    p_fb.text = "• Cost Savings: Preventative maintenance (₹1,20,000) vs catastrophic engine failure (₹35,00,000) yields ₹33,80,000+ net PLM savings.\n• PTC Windchill Export: One-click export of Service Bill of Materials (sBOM.csv) containing part numbers and live AI health state.\n• Closed-Loop Digital Thread: Connects As-Designed CAD definition with As-Operated field telemetry."
    p_fb.font.size = Pt(13)
    p_fb.font.color.rgb = WHITE_TEXT

    # ----------------------------------------------------------------------------------
    # SLIDE 10: CONCLUSION & FUTURE SCOPE
    # ----------------------------------------------------------------------------------
    slide10 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide10)
    add_header(slide10, "Conclusion & Future Scope")

    # Left: Conclusion
    c_fin = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0))
    c_fin.fill.solid()
    c_fin.fill.fore_color.rgb = CARD_BG
    c_fin.line.color.rgb = GREEN_ACCENT
    tf_cf = c_fin.text_frame
    tf_cf.word_wrap = True
    p_cfh = tf_cf.paragraphs[0]
    p_cfh.text = "✅ Project Conclusion"
    p_cfh.font.size = Pt(18)
    p_cfh.font.bold = True
    p_cfh.font.color.rgb = GREEN_ACCENT

    p_cfb = tf_cf.add_paragraph()
    p_cfb.text = "• Successfully built an AI-driven Digital Twin framework for heavy equipment engines.\n• Achieved 96.67% health classification accuracy & 0.97 R² RUL regression score.\n• Closed the loop between sensor telemetry and PLM decision making, shifting maintenance from reactive to fully predictive."
    p_cfb.font.size = Pt(14)
    p_cfb.font.color.rgb = WHITE_TEXT

    # Right: Future Scope
    f_fut = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    f_fut.fill.solid()
    f_fut.fill.fore_color.rgb = CARD_BG
    f_fut.line.color.rgb = CYAN_TITLE
    tf_ff = f_fut.text_frame
    tf_ff.word_wrap = True
    p_ffh = tf_ff.paragraphs[0]
    p_ffh.text = "🔮 Future Scope"
    p_ffh.font.size = Pt(18)
    p_ffh.font.bold = True
    p_ffh.font.color.rgb = CYAN_TITLE

    p_ffb = tf_ff.add_paragraph()
    p_ffb.text = "1. Edge AI Microcontrollers: Deploying quantized ONNX models on engine ECU edge hardware.\n2. Physics-Informed Neural Networks (PINN): Coupling thermodynamic equations with deep learning.\n3. Direct SAP / Siemens Teamcenter API Connectors: Automated work-order dispatch via MQTT/REST."
    p_ffb.font.size = Pt(14)
    p_ffb.font.color.rgb = WHITE_TEXT

    # Save presentation
    ppt_path = "PLM_Engine_Digital_Twin_Presentation.pptx"
    prs.save(ppt_path)
    print(f"Successfully generated PowerPoint presentation at: {ppt_path}")

if __name__ == "__main__":
    create_presentation()
